"""
Slide Generator Service
Creates polished McKinsey-style PowerPoint presentations using python-pptx.
Includes visual elements: icon shapes, KPI callout boxes, charts, and
decorative graphics for a professional, consulting-grade look.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def create_mckinsey_deck(deck_plan: dict, brand: dict = None, output_path: str = "output.pptx") -> str:
    """
    Generate a McKinsey-style PPTX deck from a structured deck plan.
    """
    # Defaults
    primary_color = brand.get("primary_color", "#1B3A5C") if brand else "#1B3A5C"
    accent_color = brand.get("accent_color", "#E8792F") if brand else "#E8792F"
    bg_color = brand.get("background_color", "#FFFFFF") if brand else "#FFFFFF"
    font_heading = brand.get("font_heading", "Calibri") if brand else "Calibri"
    font_body = brand.get("font_body", "Calibri") if brand else "Calibri"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = deck_plan.get("slides", [])

    for i, slide_data in enumerate(slides_data):
        layout_type = slide_data.get("layout", "title_and_content")

        if layout_type == "title":
            _add_title_slide(prs, slide_data, primary_color, accent_color, font_heading, font_body)
        elif layout_type == "section_header":
            _add_section_slide(prs, slide_data, primary_color, accent_color, font_heading)
        elif layout_type == "two_column":
            _add_two_column_slide(prs, slide_data, primary_color, accent_color, font_heading, font_body)
        elif layout_type == "process_flow":
            _add_process_flow_slide(prs, slide_data, primary_color, accent_color, font_heading, font_body)
        elif layout_type == "kpi_dashboard":
            _add_kpi_dashboard_slide(prs, slide_data, primary_color, accent_color, font_heading, font_body)
        else:
            _add_content_slide(prs, slide_data, primary_color, accent_color, font_heading, font_body)

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
    prs.save(output_path)
    return output_path


def _add_accent_bar(slide, color_hex, x=0, y=0, width=Inches(13.333), height=Inches(0.08)):
    """Add a thin accent bar to the slide."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(color_hex)
    shape.line.fill.background()


def _add_footer_bar(slide, color_hex, slide_num: int = None):
    """Add a subtle footer bar."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.1),
        Inches(13.333), Inches(0.4),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(color_hex)
    bar.line.fill.background()

    if slide_num is not None:
        tf = bar.text_frame
        tf.text = str(slide_num)
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


# ── Visual Element Helpers ──────────────────────────────────────────────────

# Icon-shape mapping: keywords → MSO_SHAPE for smart visual selection
ICON_SHAPES = {
    "cloud": MSO_SHAPE.CLOUD,
    "security": MSO_SHAPE.OCTAGON,
    "target": MSO_SHAPE.OVAL,
    "goal": MSO_SHAPE.OVAL,
    "process": MSO_SHAPE.CHEVRON,
    "step": MSO_SHAPE.CHEVRON,
    "phase": MSO_SHAPE.CHEVRON,
    "data": MSO_SHAPE.HEXAGON,
    "analytics": MSO_SHAPE.HEXAGON,
    "ai": MSO_SHAPE.DIAMOND,
    "ml": MSO_SHAPE.DIAMOND,
    "team": MSO_SHAPE.DECAGON,
    "people": MSO_SHAPE.DECAGON,
    "cost": MSO_SHAPE.DOWN_ARROW,
    "savings": MSO_SHAPE.DOWN_ARROW,
    "growth": MSO_SHAPE.UP_ARROW,
    "increase": MSO_SHAPE.UP_ARROW,
    "performance": MSO_SHAPE.LIGHTNING_BOLT,
    "speed": MSO_SHAPE.LIGHTNING_BOLT,
    "check": MSO_SHAPE.FLOWCHART_EXTRACT,
    "compliance": MSO_SHAPE.OCTAGON,
    "risk": MSO_SHAPE.FLOWCHART_EXTRACT,
    "time": MSO_SHAPE.CIRCULAR_ARROW,
    "timeline": MSO_SHAPE.RIGHT_ARROW,
    "migrate": MSO_SHAPE.RIGHT_ARROW,
    "strategy": MSO_SHAPE.STAR_5_POINT,
    "key": MSO_SHAPE.STAR_5_POINT,
    "default": MSO_SHAPE.ROUNDED_RECTANGLE,
}


def _pick_icon_shape(text: str) -> MSO_SHAPE:
    """Pick a relevant icon shape based on text keywords."""
    text_lower = text.lower()
    for keyword, shape in ICON_SHAPES.items():
        if keyword in text_lower:
            return shape
    return ICON_SHAPES["default"]


def _add_numbered_icon(slide, number: int, x, y, size, color_hex, font_size=Pt(18)):
    """Add a simple numbered badge - circle with number overlay."""
    # Create circle background
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = hex_to_rgb(color_hex)
    circle.line.fill.background()
    
    # Create a text box on top of circle for the number
    txt_box = slide.shapes.add_textbox(x, y, size, size)
    tf = txt_box.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    run = tf.paragraphs[0].add_run()
    run.text = str(number)
    run.font.name = "Arial Black"
    run.font.size = font_size
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    
    return circle


def _add_icon_shape(slide, shape_type, x, y, width, height, fill_color_hex, icon_text=""):
    """Add a decorative icon shape with optional text inside."""
    shape = slide.shapes.add_shape(shape_type, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color_hex)
    shape.line.fill.background()
    if icon_text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = icon_text
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    return shape


def _add_kpi_card(slide, x, y, width, height, value, label, color_hex, accent_hex):
    """Add a KPI metric card with large number and full label (no truncation)."""
    import re
    
    # Clean up value - keep it SHORT for display
    value_str = str(value).strip()
    # Extract just the number part
    num_match = re.match(r'^([\$]?[\d,]+\.?\d*[%MBKx+]?\+?)', value_str)
    if num_match:
        value_str = num_match.group(1)
    
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(245, 247, 250)
    card.line.color.rgb = RGBColor(210, 215, 225)
    card.line.width = Pt(1)

    # Accent strip on top of card
    accent_strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x + Inches(0.03), y + Inches(0.03),
        width - Inches(0.06), Inches(0.08)
    )
    accent_strip.fill.solid()
    accent_strip.fill.fore_color.rgb = hex_to_rgb(accent_hex)
    accent_strip.line.fill.background()

    # Value text - auto-size font based on length
    val_len = len(value_str)
    if val_len <= 4:
        font_size = Pt(42)
    elif val_len <= 6:
        font_size = Pt(36)
    elif val_len <= 8:
        font_size = Pt(30)
    else:
        font_size = Pt(24)
    
    val_box = slide.shapes.add_textbox(x, y + Inches(0.2), width, Inches(1.0))
    tf = val_box.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = value_str
    p.font.name = "Arial Black"
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(color_hex)
    p.alignment = PP_ALIGN.CENTER

    # Label text - FULL height available, auto-size font, NO TRUNCATION
    label_clean = label.strip()
    # Choose font size based on label length
    if len(label_clean) <= 25:
        label_font = Pt(11)
    elif len(label_clean) <= 40:
        label_font = Pt(10)
    elif len(label_clean) <= 55:
        label_font = Pt(9)
    else:
        label_font = Pt(8)
    
    # Label takes up bottom portion of card - make it tall enough
    label_height = height - Inches(1.3)
    lbl_box = slide.shapes.add_textbox(x + Inches(0.05), y + Inches(1.2), width - Inches(0.1), label_height)
    tf2 = lbl_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = label_clean  # No truncation - show full text
    p2.font.name = "Arial"
    p2.font.size = label_font
    p2.font.color.rgb = RGBColor(80, 90, 100)
    p2.alignment = PP_ALIGN.CENTER


def _add_process_chevrons(slide, items, y, accent_hex, primary_hex):
    """Add a horizontal process flow with chevron arrows."""
    num = min(len(items), 5)
    if num == 0:
        return
    total_width = 11.0
    chevron_w = total_width / num
    start_x = 1.2

    for i, item in enumerate(items[:5]):
        x = start_x + i * chevron_w
        # Alternate colors
        fill = accent_hex if i % 2 == 0 else primary_hex
        chev = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            Inches(x), Inches(y),
            Inches(chevron_w - 0.1), Inches(1.0)
        )
        chev.fill.solid()
        chev.fill.fore_color.rgb = hex_to_rgb(fill)
        chev.line.fill.background()
        tf = chev.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        # Truncate text for chevrons
        p.text = item[:30] if len(item) > 30 else item
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER


def _extract_kpis(bullet_points: list) -> list:
    """Extract numeric KPI values from bullet points. Keep full labels - no truncation."""
    import re
    kpis = []
    for bp in bullet_points:
        # Match numeric patterns: "35%", "$4.2M", "99.99%", "50+", "18", "4.5/5.0"
        match = re.search(r'([\$]?[\d,]+\.?\d*[%MBK]?\+?(?:/[\d.]+)?)', bp)
        
        if match:
            value = match.group(1).strip()
            
            # Keep value SHORT - just the number
            value = re.sub(r'\s*(applications?|apps?|engineers?|months?|years?|users?|days?).*$', '', value, flags=re.IGNORECASE)
            
            # Skip if value is just a single digit that might be a list number
            if re.match(r'^\d$', value) and bp.strip().startswith(value):
                continue
            
            # Extract label from the rest of the text
            label = bp.replace(match.group(0), "").strip(" -:–—•▸,")
            
            # Clean up label - remove leading/trailing punctuation
            label = re.sub(r'^[\s\-:–—•▸,]+', '', label)
            label = re.sub(r'[\s\-:–—•▸,]+$', '', label)
            
            # Capitalize first letter
            if label:
                label = label[0].upper() + label[1:] if len(label) > 1 else label.upper()
            
            # NO TRUNCATION - keep full label, the card will handle it
            if not label:
                label = "Key Metric"
            
            kpis.append({"value": value, "label": label})
    return kpis


def _add_title_slide(prs, data, primary, accent, font_h, font_b):
    """Create a polished title slide with full-width color block and decorative shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Full slide dark background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = hex_to_rgb(primary)
    bg_shape.line.fill.background()

    # Accent bar at top
    _add_accent_bar(slide, accent, y=Inches(0), height=Inches(0.12))

    # Decorative geometric shapes (semi-transparent circles)
    for cx, cy, sz, opacity in [
        (10.5, 0.5, 2.5, 0.08),
        (11.5, 5.0, 3.0, 0.06),
        (0.2, 6.0, 1.5, 0.07),
        (9.0, 2.5, 1.0, 0.10),
    ]:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(sz), Inches(sz)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = hex_to_rgb(accent)
        circle.fill.fore_color.brightness = 0.4
        circle.line.fill.background()

    # Decorative diamond
    diamond = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND, Inches(12.0), Inches(3.5), Inches(0.8), Inches(0.8)
    )
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = hex_to_rgb(accent)
    diamond.line.fill.background()

    # Title - adjusted position to avoid overlap
    title_text = data.get("title", "Executive Presentation")
    # Use smaller font for long titles
    title_font_size = Pt(38) if len(title_text) > 40 else Pt(44)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(9), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = font_h
    p.font.size = title_font_size
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    # Decorative line - position it after title area
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(4.5),
        Inches(3.5), Inches(0.05),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(accent)
    line.line.fill.background()

    # Subtitle - positioned below the line
    subtitle = data.get("subtitle", "")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.75), Inches(9), Inches(1))
        tf2 = sub_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.name = font_b
        p2.font.size = Pt(20)
        p2.font.color.rgb = hex_to_rgb(accent)
        p2.alignment = PP_ALIGN.LEFT

    # "Generated by SlideForge AI" watermark
    wm = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(4), Inches(0.5))
    wm_tf = wm.text_frame
    p_wm = wm_tf.paragraphs[0]
    p_wm.text = "Generated by SlideForge AI"
    p_wm.font.size = Pt(10)
    p_wm.font.color.rgb = RGBColor(180, 190, 210)
    p_wm.font.italic = True


def _add_section_slide(prs, data, primary, accent, font_h):
    """Create a section divider slide with decorative shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Left color block - wider to accommodate text
    left_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(6), Inches(7.5),
    )
    left_block.fill.solid()
    left_block.fill.fore_color.rgb = hex_to_rgb(primary)
    left_block.line.fill.background()

    # Decorative circles on left block
    for cx, cy, sz in [(4.5, 0.5, 1.5), (0.3, 5.5, 1.0), (5.0, 6.2, 0.6)]:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(sz), Inches(sz))
        c.fill.solid()
        c.fill.fore_color.rgb = hex_to_rgb(accent)
        c.fill.fore_color.brightness = 0.35
        c.line.fill.background()

    # Section title on left block - larger area, smaller font for long titles
    title_text = data.get("title", "")
    # Use smaller font if title is long
    title_font_size = Pt(32) if len(title_text) > 25 else Pt(38)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(4.8), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = font_h
    p.font.size = title_font_size
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    # Accent bar
    _add_accent_bar(slide, accent, x=Inches(0.8), y=Inches(5.8), width=Inches(2.5), height=Inches(0.06))

    # Right side content
    if data.get("bullet_points"):
        content_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5))
        tf2 = content_box.text_frame
        tf2.word_wrap = True
        for j, point in enumerate(data["bullet_points"][:4]):  # Max 4 points
            p2 = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()

            # Add numbered icon for each point
            _add_numbered_icon(
                slide, j + 1,
                Inches(6.4), Inches(2.1 + j * 1.0),
                Inches(0.35), accent, font_size=Pt(12)
            )

            p2.text = f"    {point}"
            p2.font.name = "Calibri"
            p2.font.size = Pt(15)
            p2.font.color.rgb = RGBColor(60, 60, 60)
            p2.space_after = Pt(20)

    _add_footer_bar(slide, primary)


def _add_content_slide(prs, data, primary, accent, font_h, font_b):
    """Create a clean content slide with numbered bullets - no KPI cards here."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent bar
    _add_accent_bar(slide, accent, y=Inches(0), height=Inches(0.06))

    # Headline bar (dark background for the headline)
    headline_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.06),
        Inches(13.333), Inches(1.2),
    )
    headline_bg.fill.solid()
    headline_bg.fill.fore_color.rgb = hex_to_rgb(primary)
    headline_bg.line.fill.background()

    # Headline text
    headline_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(1.0))
    tf = headline_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = data.get("title", "")
    p.font.name = font_h
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    # Content area - always use simple numbered bullets for reliability
    bullets = data.get("bullet_points", [])
    
    if bullets:
        y_start = 1.7
        for j, point in enumerate(bullets):
            y_pos = y_start + j * 0.85
            if y_pos > 6.2:
                break

            # Number badge
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.6), Inches(y_pos),
                Inches(0.45), Inches(0.45)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = hex_to_rgb(accent)
            badge.line.fill.background()
            
            # Number text overlay
            num_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(y_pos),
                Inches(0.45), Inches(0.45)
            )
            num_tf = num_box.text_frame
            num_tf.word_wrap = False
            num_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            num_p = num_tf.paragraphs[0]
            num_p.text = str(j + 1)
            num_p.font.name = "Arial Black"
            num_p.font.size = Pt(14)
            num_p.font.bold = True
            num_p.font.color.rgb = RGBColor(255, 255, 255)
            num_p.alignment = PP_ALIGN.CENTER

            # Bullet text - full width
            txt_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(y_pos),
                Inches(11.5), Inches(0.75)
            )
            tf2 = txt_box.text_frame
            tf2.word_wrap = True
            tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
            p2 = tf2.paragraphs[0]
            p2.text = point
            p2.font.name = "Arial"
            p2.font.size = Pt(16)
            p2.font.color.rgb = RGBColor(50, 50, 50)

    # Left accent side-bar decoration
    side_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.26),
        Inches(0.08), Inches(5.84),
    )
    side_bar.fill.solid()
    side_bar.fill.fore_color.rgb = hex_to_rgb(accent)
    side_bar.line.fill.background()

    # Speaker notes
    notes = data.get("speaker_notes", "")
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    _add_footer_bar(slide, primary)


def _add_two_column_slide(prs, data, primary, accent, font_h, font_b):
    """Create a two-column content slide with icon bullets and visual divider."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent bar
    _add_accent_bar(slide, accent, y=Inches(0), height=Inches(0.06))

    # Headline bar
    headline_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.06),
        Inches(13.333), Inches(1.2),
    )
    headline_bg.fill.solid()
    headline_bg.fill.fore_color.rgb = hex_to_rgb(primary)
    headline_bg.line.fill.background()

    headline_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(1.0))
    tf = headline_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = data.get("title", "")
    p.font.name = font_h
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    # Split bullet points into two columns
    bullets = data.get("bullet_points", [])
    mid = len(bullets) // 2 if len(bullets) > 1 else len(bullets)
    left_bullets = bullets[:mid]
    right_bullets = bullets[mid:]

    # Left column with numbered icons
    y_start = 1.8
    for j, point in enumerate(left_bullets):
        y_pos = y_start + j * 0.8
        if y_pos > 6.0:
            break
        _add_numbered_icon(slide, j + 1, Inches(0.8), Inches(y_pos), Inches(0.45), accent, Pt(14))
        txt = slide.shapes.add_textbox(Inches(1.45), Inches(y_pos + 0.02), Inches(4.8), Inches(0.65))
        tf_l = txt.text_frame
        tf_l.word_wrap = True
        tf_l.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_l = tf_l.paragraphs[0]
        p_l.text = point
        p_l.font.name = font_b
        p_l.font.size = Pt(15)
        p_l.font.color.rgb = RGBColor(50, 50, 50)

    # Decorative vertical divider with accent dot
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.6), Inches(1.8),
        Inches(0.03), Inches(4.5),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = hex_to_rgb(accent)
    divider.line.fill.background()

    # Diamond at center of divider
    diamond = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        Inches(6.45), Inches(3.9),
        Inches(0.35), Inches(0.35),
    )
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = hex_to_rgb(accent)
    diamond.line.fill.background()

    # Right column with numbered icons
    for j, point in enumerate(right_bullets):
        y_pos = y_start + j * 0.8
        if y_pos > 6.0:
            break
        _add_numbered_icon(slide, mid + j + 1, Inches(7.0), Inches(y_pos), Inches(0.45), primary, Pt(14))
        txt = slide.shapes.add_textbox(Inches(7.65), Inches(y_pos + 0.02), Inches(4.8), Inches(0.65))
        tf_r = txt.text_frame
        tf_r.word_wrap = True
        tf_r.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_r = tf_r.paragraphs[0]
        p_r.text = point
        p_r.font.name = font_b
        p_r.font.size = Pt(15)
        p_r.font.color.rgb = RGBColor(50, 50, 50)

    # Left accent side-bar
    side_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.26),
        Inches(0.08), Inches(5.84),
    )
    side_bar.fill.solid()
    side_bar.fill.fore_color.rgb = hex_to_rgb(accent)
    side_bar.line.fill.background()

    # Speaker notes
    notes = data.get("speaker_notes", "")
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    _add_footer_bar(slide, primary)


# ── Additional visual layout types ──────────────────────────────────────────


def _add_process_flow_slide(prs, data, primary, accent, font_h, font_b):
    """Create a process/timeline slide with clean step cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_accent_bar(slide, accent, y=Inches(0), height=Inches(0.06))

    headline_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.06),
        Inches(13.333), Inches(1.2),
    )
    headline_bg.fill.solid()
    headline_bg.fill.fore_color.rgb = hex_to_rgb(primary)
    headline_bg.line.fill.background()

    headline_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(1.0))
    tf = headline_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = data.get("title", "")
    p.font.name = font_h
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    bullets = data.get("bullet_points", [])
    num = min(len(bullets), 5)
    
    if num > 0:
        total_width = 12.0
        box_w = total_width / num
        gap = 0.15
        
        for i, point in enumerate(bullets[:5]):
            x = 0.65 + i * box_w
            fill_color = accent if i % 2 == 0 else primary
            
            # Card container
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.6),
                Inches(box_w - gap), Inches(5.0),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(250, 251, 252)
            card.line.color.rgb = RGBColor(210, 218, 225)
            card.line.width = Pt(1)

            # Top accent strip
            card_accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(1.6),
                Inches(box_w - gap), Inches(0.1),
            )
            card_accent.fill.solid()
            card_accent.fill.fore_color.rgb = hex_to_rgb(fill_color)
            card_accent.line.fill.background()

            # Step number badge - simple colored square with number
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + (box_w - gap) / 2 - 0.25), Inches(1.85),
                Inches(0.5), Inches(0.5),
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = hex_to_rgb(fill_color)
            badge.line.fill.background()
            
            # Number text on badge
            num_box = slide.shapes.add_textbox(
                Inches(x + (box_w - gap) / 2 - 0.25), Inches(1.85),
                Inches(0.5), Inches(0.5)
            )
            num_tf = num_box.text_frame
            num_tf.word_wrap = False
            num_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            num_p = num_tf.paragraphs[0]
            num_p.text = str(i + 1)
            num_p.font.name = "Arial Black"
            num_p.font.size = Pt(18)
            num_p.font.bold = True
            num_p.font.color.rgb = RGBColor(255, 255, 255)
            num_p.alignment = PP_ALIGN.CENTER
            
            # Parse phase name and description
            if ':' in point:
                parts = point.split(':', 1)
                phase_name = parts[0].strip()
                description = parts[1].strip() if len(parts) > 1 else ""
            else:
                words = point.split()
                if len(words) > 3:
                    phase_name = " ".join(words[:2])
                    description = " ".join(words[2:])
                else:
                    phase_name = point
                    description = ""

            # Phase name header
            phase_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(2.5),
                Inches(box_w - gap - 0.2), Inches(0.7),
            )
            tf_phase = phase_box.text_frame
            tf_phase.word_wrap = True
            p_phase = tf_phase.paragraphs[0]
            p_phase.text = phase_name
            p_phase.font.name = "Arial"
            p_phase.font.size = Pt(12)
            p_phase.font.bold = True
            p_phase.font.color.rgb = hex_to_rgb(primary)
            p_phase.alignment = PP_ALIGN.CENTER

            # Description text - generous area
            if description:
                desc_box = slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(3.25),
                    Inches(box_w - gap - 0.2), Inches(3.1),
                )
                tf_desc = desc_box.text_frame
                tf_desc.word_wrap = True
                p_desc = tf_desc.paragraphs[0]
                p_desc.text = description
                p_desc.font.name = "Arial"
                p_desc.font.size = Pt(10)
                p_desc.font.color.rgb = RGBColor(60, 65, 70)
                p_desc.alignment = PP_ALIGN.CENTER

            # Arrow connector (except last)
            if i < num - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Inches(x + box_w - gap - 0.05), Inches(3.8),
                    Inches(0.25), Inches(0.2),
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = hex_to_rgb(accent)
                arrow.line.fill.background()

    # Left accent side-bar
    side_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.26),
        Inches(0.08), Inches(5.84),
    )
    side_bar.fill.solid()
    side_bar.fill.fore_color.rgb = hex_to_rgb(accent)
    side_bar.line.fill.background()

    notes = data.get("speaker_notes", "")
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    _add_footer_bar(slide, primary)


def _add_kpi_dashboard_slide(prs, data, primary, accent, font_h, font_b):
    """Create a KPI dashboard slide with large metric cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_accent_bar(slide, accent, y=Inches(0), height=Inches(0.06))

    headline_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.06),
        Inches(13.333), Inches(1.2),
    )
    headline_bg.fill.solid()
    headline_bg.fill.fore_color.rgb = hex_to_rgb(primary)
    headline_bg.line.fill.background()

    headline_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(1.0))
    tf = headline_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = data.get("title", "")
    p.font.name = font_h
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    bullets = data.get("bullet_points", [])
    kpis = _extract_kpis(bullets)

    if kpis:
        num_kpis = len(kpis)
        
        if num_kpis <= 3:
            # Single row - larger cards
            card_w = 3.2
            card_h = 2.4
            total_w = num_kpis * card_w + (num_kpis - 1) * 0.4
            start_x = (13.333 - total_w) / 2
            
            for idx, kpi in enumerate(kpis):
                x = start_x + idx * (card_w + 0.4)
                _add_kpi_card(
                    slide, Inches(x), Inches(2.5),
                    Inches(card_w), Inches(card_h),
                    kpi["value"], kpi["label"],
                    primary, accent
                )
        elif num_kpis <= 4:
            # Single row of 4 - medium cards
            card_w = 2.8
            card_h = 2.3
            gap = 0.35
            total_w = num_kpis * card_w + (num_kpis - 1) * gap
            start_x = (13.333 - total_w) / 2
            
            for idx, kpi in enumerate(kpis):
                x = start_x + idx * (card_w + gap)
                _add_kpi_card(
                    slide, Inches(x), Inches(2.3),
                    Inches(card_w), Inches(card_h),
                    kpi["value"], kpi["label"],
                    primary, accent
                )
        else:
            # Two rows
            card_w = 2.7
            card_h = 2.1
            gap = 0.3
            
            # Top row (first 4)
            top_count = min(4, num_kpis)
            total_w = top_count * card_w + (top_count - 1) * gap
            start_x = (13.333 - total_w) / 2
            
            for idx, kpi in enumerate(kpis[:4]):
                x = start_x + idx * (card_w + gap)
                _add_kpi_card(
                    slide, Inches(x), Inches(1.7),
                    Inches(card_w), Inches(card_h),
                    kpi["value"], kpi["label"],
                    primary, accent
                )
            
            # Bottom row (remaining)
            if num_kpis > 4:
                bottom_count = min(4, num_kpis - 4)
                total_w = bottom_count * card_w + (bottom_count - 1) * gap
                start_x = (13.333 - total_w) / 2
                
                for idx, kpi in enumerate(kpis[4:8]):
                    x = start_x + idx * (card_w + gap)
                    _add_kpi_card(
                        slide, Inches(x), Inches(4.1),
                        Inches(card_w), Inches(card_h),
                        kpi["value"], kpi["label"],
                        primary, accent
                    )
    else:
        # Fallback to bullet list if no KPIs extracted
        for j, point in enumerate(bullets):
            y_pos = 1.8 + j * 0.9
            if y_pos > 6.2:
                break
            
            # Simple bullet marker
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(0.8), Inches(y_pos + 0.15),
                Inches(0.15), Inches(0.15)
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = hex_to_rgb(accent)
            marker.line.fill.background()
            
            txt = slide.shapes.add_textbox(Inches(1.15), Inches(y_pos), Inches(11.0), Inches(0.8))
            tf2 = txt.text_frame
            tf2.word_wrap = True
            tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
            p2 = tf2.paragraphs[0]
            p2.text = point
            p2.font.name = "Arial"
            p2.font.size = Pt(16)
            p2.font.color.rgb = RGBColor(50, 50, 50)

    notes = data.get("speaker_notes", "")
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    _add_footer_bar(slide, primary)
